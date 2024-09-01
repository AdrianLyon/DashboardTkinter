import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageTk
import io

class PPTXViewerApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PPTX Viewer")
        self.geometry("800x600")
        self.configure(bg="#f0f0f0")

        self.prs = None
        self.current_slide = 0

        self.create_widgets()

    def create_widgets(self):
        self.btn_open = ttk.Button(self, text="Open PPTX File", command=self.open_pptx)
        self.btn_open.pack(pady=20)

        self.canvas = tk.Canvas(self, bg="white", width=640, height=480)
        self.canvas.pack(pady=20)

        self.btn_change_image = ttk.Button(self, text="Change Image", command=self.change_image)
        self.btn_change_image.pack(pady=10)

    def open_pptx(self):
        file_path = filedialog.askopenfilename(
            title="Select PPTX File",
            filetypes=(("PowerPoint Files", "*.pptx"), ("All Files", "*.*"))
        )
        if file_path:
            self.prs = Presentation(file_path)
            self.current_slide = 0
            self.show_slide(self.current_slide)

    def show_slide(self, slide_index):
        if not self.prs:
            return

        slide = self.prs.slides[slide_index]
        self.canvas.delete("all")

        y_position = 20
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                self.canvas.create_text(320, y_position, text=text, fill="black", font=("Arial", 14))
                y_position += 40

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = shape.image.blob
                try:
                    image = Image.open(io.BytesIO(image_stream))
                    image.thumbnail((640, 480), Image.Resampling.LANCZOS)
                    img = ImageTk.PhotoImage(image)
                    self.canvas.create_image(320, 240, image=img)
                    self.canvas.image = img
                except OSError:
                    continue  # Skip unsupported images

    def change_image(self):
        if not self.prs:
            messagebox.showerror("Error", "No PPTX file is loaded.")
            return

        slide = self.prs.slides[self.current_slide]

        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pictures:
            messagebox.showinfo("No Pictures", "No pictures found in this slide.")
            return

        file_path = filedialog.askopenfilename(
            title="Select New Image",
            filetypes=(("Image Files", "*.jpg;*.jpeg;*.png;*.bmp"), ("All Files", "*.*"))
        )
        if file_path:
            for picture in pictures:
                picture.image = self.prs.part.related_parts[picture._element.blip_rId].blob = open(file_path, 'rb').read()
            messagebox.showinfo("Success", "Image replaced successfully.")
            self.show_slide(self.current_slide)

if __name__ == "__main__":
    app = PPTXViewerApp()
    app.mainloop()
