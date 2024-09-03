import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from pptx import Presentation
from pptx_replace import replace_text,replace_picture

class PresentationGeneratorApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Presentation Generator")
        self.geometry("600x600")
        self.configure(bg="#f0f0f0")

        self.create_widgets()

    def create_widgets(self):
        # Labels and Entry fields
        self.fields = {
            "nombre": tk.StringVar(),
            "curp": tk.StringVar(),
            "categoria": tk.StringVar(),
            "curso": tk.StringVar(),
            "folio": tk.StringVar(),
            "imagen_nueva": tk.StringVar()
        }

        row = 0
        for label, var in self.fields.items():
            tk.Label(self, text=label, bg="#f0f0f0").grid(row=row, column=0, sticky=tk.W, padx=10, pady=5)
            tk.Entry(self, textvariable=var, width=30).grid(row=row, column=1, padx=10, pady=5)
            row += 1

        # Button to open pptx template
        self.btn_open = ttk.Button(self, text="Select PPTX Template", command=self.open_template)
        self.btn_open.grid(row=row, column=0, padx=10, pady=20, columnspan=2)

        # Button to generate pptx
        self.btn_generate = ttk.Button(self, text="Generate Presentation", command=self.generate_presentation)
        self.btn_generate.grid(row=row+1, column=0, padx=10, pady=10, columnspan=2)
        
        # Botón para seleccionar la imagen
        self.btn_image = ttk.Button(self, text="Select Image", command=self.select_image)
        self.btn_image.grid(row=row, column=0, padx=10, pady=5, columnspan=2)

        self.ppt_template = None

    def open_template(self):
        file_path = filedialog.askopenfilename(
            title="Select PPTX Template",
            filetypes=(("PowerPoint Files", "*.pptx"), ("All Files", "*.*"))
        )
        if file_path:
            self.ppt_template = file_path
            messagebox.showinfo("Template Selected", f"Selected Template: {file_path}")
        else:
            messagebox.showwarning("No File Selected", "Please select a PPTX template to proceed.")

    def select_image(self):
        file_path = filedialog.askopenfilename(
            title="Select Image",
            filetypes=(("Image Files", "*.png;*.jpg;*.jpeg;*.bmp;*.gif"), ("All Files", "*.*"))
        )
        if file_path:
            self.fields["imagen_nueva"].set(file_path)
            messagebox.showinfo("Image Selected", f"Selected Image: {file_path}")
        else:
            messagebox.showwarning("No File Selected", "Please select an image to proceed.")

    def generate_presentation(self):
        if not self.ppt_template:
            messagebox.showerror("Error", "Please select a PPTX template before generating the presentation.")
            return

        # Generar el contexto, asegurándonos de que todas las entradas sean cadenas de texto
        context = {key.lower().replace(" ", "_"): str(var.get()) for key, var in self.fields.items()}
        
        try:
            prs = Presentation(self.ppt_template)

            # Reemplazar el texto en las diapositivas
            for key, value in context.items():
                replace_text(prs, f"{{{key}}}", value)

            # Reemplazar la imagen si se ha seleccionado una
            image_path = self.fields["imagen_nueva"].get()
            if image_path:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if "{imagen_nueva}" in run.text:
                                        run.text = ""  # Eliminar el marcador
                                        # Insertar la imagen en la posición del marcador
                                        slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                                        break

            # Guardar la presentación modificada
            output_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=(("PowerPoint Files", "*.pptx"), ("All Files", "*.*")),
                title="Save Generated Presentation"
            )
            if output_path:
                prs.save(output_path)
                messagebox.showinfo("Success", f"Presentation saved successfully: {output_path}")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")




if __name__ == "__main__":
    app = PresentationGeneratorApp()
    app.mainloop()
