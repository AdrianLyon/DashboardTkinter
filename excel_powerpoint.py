import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
from pptx import Presentation
import pandas as pd

class BatchPresentationGeneratorApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Batch Presentation Generator")
        self.geometry("600x300")
        self.configure(bg="#f0f0f0")

        self.create_widgets()

    def create_widgets(self):
        # Buttons to open pptx template and excel file
        self.btn_open_pptx = ttk.Button(self, text="Select PPTX Template", command=self.open_pptx_template)
        self.btn_open_pptx.grid(row=0, column=0, padx=10, pady=20, columnspan=2)

        self.btn_open_excel = ttk.Button(self, text="Select Excel File", command=self.open_excel_file)
        self.btn_open_excel.grid(row=1, column=0, padx=10, pady=10, columnspan=2)

        # Button to generate presentations
        self.btn_generate = ttk.Button(self, text="Generate Presentations", command=self.generate_presentations)
        self.btn_generate.grid(row=2, column=0, padx=10, pady=10, columnspan=2)

        self.pptx_template = None
        self.excel_file = None

    def open_pptx_template(self):
        file_path = filedialog.askopenfilename(
            title="Select PPTX Template",
            filetypes=(("PowerPoint Files", "*.pptx"), ("All Files", "*.*"))
        )
        if file_path:
            self.pptx_template = file_path
            messagebox.showinfo("Template Selected", f"Selected Template: {file_path}")
        else:
            messagebox.showwarning("No File Selected", "Please select a PPTX template to proceed.")

    def open_excel_file(self):
        file_path = filedialog.askopenfilename(
            title="Select Excel File",
            filetypes=(("Excel Files", "*.xlsx"), ("All Files", "*.*"))
        )
        if file_path:
            self.excel_file = file_path
            messagebox.showinfo("Excel File Selected", f"Selected Excel File: {file_path}")
        else:
            messagebox.showwarning("No File Selected", "Please select an Excel file to proceed.")

    def generate_presentations(self):
        if not self.pptx_template:
            messagebox.showerror("Error", "Please select a PPTX template before generating the presentations.")
            return
        if not self.excel_file:
            messagebox.showerror("Error", "Please select an Excel file before generating the presentations.")
            return

        try:
            df = pd.read_excel(self.excel_file)
            for index, row in df.iterrows():
                prs = Presentation(self.pptx_template)
                
                context = {
                    'nombre': row['nombre'],
                    'curp': row['curp'],
                    'categoria': row['categoria'],
                    'curso': row['curso']
                }

                # Replace text in slides with context data
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            full_text = "".join([run.text for run in paragraph.runs])
                            for key, value in context.items():
                                if f"{{{key}}}" in full_text:
                                    full_text = full_text.replace(f"{{{key}}}", value)

                            # Update each run with the new text
                            current_pos = 0
                            for run in paragraph.runs:
                                run_len = len(run.text)
                                run.text = full_text[current_pos:current_pos + run_len]
                                current_pos += run_len

                output_path = f"generated_presentation_{index}.pptx"
                prs.save(output_path)

            messagebox.showinfo("Success", "Presentations generated successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")

if __name__ == "__main__":
    app = BatchPresentationGeneratorApp()
    app.mainloop()
